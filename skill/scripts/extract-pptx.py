#!/usr/bin/env python3
"""
Enhanced PPTX content extractor for ppt-to-course skill.

Extracts ALL content from a PowerPoint file (.pptx) and produces a
slide-profile.json conforming to slide-profile.schema.json, plus
extracted images in assets/ and reference PNGs in references/.

Strategy: Text-first with OCR fallback.
- python-pptx provides structured text, formatting, tables, charts
- For image-heavy slides with <20 chars extractable text, fall back
  to rendering the slide as an image + OCR via pytesseract

Usage:
    python extract-pptx.py <input.pptx> [output_dir]

Dependencies:
    Required: python-pptx, Pillow
    Optional: pytesseract (for OCR fallback)
"""

import json
import logging
import os
import re
import struct
import sys
import uuid
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import pytesseract

    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("extract-pptx")

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
MONOSPACE_FONTS = {
    "consolas",
    "courier",
    "courier new",
    "jetbrains mono",
    "fira code",
    "fira mono",
    "source code pro",
    "cascadia code",
    "cascadia mono",
    "menlo",
    "monaco",
    "sf mono",
    "roboto mono",
    "ubuntu mono",
    "dejavu sans mono",
    "liberation mono",
    "lucida console",
    "andale mono",
}

BULLET_CHARS = {
    "\u2022",  # bullet
    "\u2023",  # triangular bullet
    "\u25aa",  # black small square
    "\u25ab",  # white small square
    "\u25cf",  # black circle
    "\u25cb",  # white circle
    "\u25a0",  # black square
    "\u25a1",  # white square
    "\u2013",  # en dash
    "\u2014",  # em dash
    "\u25b6",  # black right-pointing triangle
    "\u25b8",  # black right-pointing small triangle
    "\u2043",  # hyphen bullet
    "\u27a4",  # black right arrowhead
}


# ---------------------------------------------------------------------------
# Helper: EMU → percentage of slide dimensions
# ---------------------------------------------------------------------------
def _emu_to_pct(emu_value, slide_dimension_emu):
    """Convert EMU value to percentage of slide dimension."""
    if slide_dimension_emu == 0:
        return 0.0
    return round((emu_value / slide_dimension_emu) * 100, 2)


# ---------------------------------------------------------------------------
# Helper: Extract color from pptx run
# ---------------------------------------------------------------------------
def _extract_color(run):
    """Extract hex color from a run, if available."""
    try:
        if run.font.color and run.font.color.rgb:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# Helper: Detect bullet type from paragraph XML
# ---------------------------------------------------------------------------
def _detect_bullet_info(paragraph):
    """
    Inspect paragraph._pPr XML to detect bullet/numbering info.
    Returns ('bullet', level) | ('numbered', level) | (None, level).
    """
    level = paragraph.level or 0
    pPr = paragraph._pPr
    if pPr is None:
        return None, level

    nsmap = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }

    # Check for numbered list marker
    if pPr.find("a:buAutoNum", nsmap) is not None:
        return "numbered", level

    # Check for character/font bullet
    if pPr.find("a:buChar", nsmap) is not None:
        return "bullet", level
    if pPr.find("a:buFont", nsmap) is not None:
        return "bullet", level

    # Check for picture bullet
    if pPr.find("a:buBlip", nsmap) is not None:
        return "bullet", level

    # Check if text starts with common bullet chars
    text = paragraph.text.strip()
    if text and text[0] in BULLET_CHARS:
        return "bullet", level

    return None, level


# ---------------------------------------------------------------------------
# Helper: Classify a text block
# ---------------------------------------------------------------------------
def _classify_block(
    paragraph,
    is_title_shape,
    shape_idx,
    total_shapes,
    avg_font_size,
    all_monospace,
):
    """
    Classify a paragraph/block into one of:
    heading, subheading, body_text, bullet_list, numbered_list,
    code_block, quote, caption, formula
    """
    bullet_type, level = _detect_bullet_info(paragraph)

    # Code detection: monospace font
    if all_monospace:
        return "code_block"

    # Title shape → heading
    if is_title_shape and shape_idx == 0:
        return "heading"

    # Numbered list
    if bullet_type == "numbered":
        return "numbered_list"

    # Bullet list
    if bullet_type == "bullet":
        return "bullet_list"

    # Level > 0 without explicit bullet = indented bullet
    if level > 0:
        return "bullet_list"

    # Heading heuristics
    if avg_font_size and avg_font_size >= 24:
        return "heading"
    if avg_font_size and avg_font_size >= 18:
        return "subheading"

    # Caption heuristic: small text
    if avg_font_size and avg_font_size <= 10:
        return "caption"

    # Quote heuristic: starts with " or italicised
    text = paragraph.text.strip()
    if text.startswith('"') or text.startswith("\u201c"):
        return "quote"

    # Formula heuristic: math-heavy characters
    math_chars = set("=+-*/^{}[]()\\|_<>")
    if len(text) > 3 and sum(1 for c in text if c in math_chars) / len(text) > 0.25:
        return "formula"

    return "body_text"


# ---------------------------------------------------------------------------
# Helper: Extract paragraphs with formatting from a text frame
# ---------------------------------------------------------------------------
def _extract_paragraphs(
    text_frame,
    is_title_shape,
    shape_idx,
    total_shapes,
    slide_width_emu,
    slide_height_emu,
    shape,
):
    """Extract paragraphs from a text frame with formatting metadata."""
    blocks = []
    # Group consecutive paragraphs of same list type
    current_list_items = []
    current_list_type = None

    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            # Flush pending list
            if current_list_items:
                blocks.append(
                    _make_list_block(
                        current_list_type,
                        current_list_items,
                        shape,
                        slide_width_emu,
                        slide_height_emu,
                    )
                )
                current_list_items = []
                current_list_type = None
            continue

        # Gather formatting from runs
        runs_bold = []
        runs_italic = []
        runs_font_size = []
        runs_color = []
        runs_font_name = []
        for run in para.runs:
            if run.font.bold:
                runs_bold.append(True)
            if run.font.italic:
                runs_italic.append(True)
            if run.font.size:
                runs_font_size.append(run.font.size.pt)
            color = _extract_color(run)
            if color:
                runs_color.append(color)
            if run.font.name:
                runs_font_name.append(run.font.name.lower())

        avg_font_size = (
            sum(runs_font_size) / len(runs_font_size) if runs_font_size else None
        )
        all_monospace = (
            bool(runs_font_name)
            and all(fn in MONOSPACE_FONTS for fn in runs_font_name)
        )
        is_bold = len(runs_bold) > 0 and len(runs_bold) >= len(para.runs) * 0.5
        is_italic = len(runs_italic) > 0 and len(runs_italic) >= len(para.runs) * 0.5

        block_type = _classify_block(
            para,
            is_title_shape,
            shape_idx,
            total_shapes,
            avg_font_size,
            all_monospace,
        )

        formatting = {}
        if is_bold:
            formatting["bold"] = True
        if is_italic:
            formatting["italic"] = True
        if avg_font_size:
            formatting["font_size_pt"] = round(avg_font_size, 1)
        if runs_color:
            formatting["color"] = runs_color[0]

        # Handle list grouping
        if block_type in ("bullet_list", "numbered_list"):
            if current_list_type == block_type:
                current_list_items.append(
                    _clean_bullet_text(text) if block_type == "bullet_list" else text
                )
            else:
                # Flush previous list
                if current_list_items:
                    blocks.append(
                        _make_list_block(
                            current_list_type,
                            current_list_items,
                            shape,
                            slide_width_emu,
                            slide_height_emu,
                        )
                    )
                current_list_type = block_type
                current_list_items = [
                    _clean_bullet_text(text) if block_type == "bullet_list" else text
                ]
        else:
            # Flush pending list
            if current_list_items:
                blocks.append(
                    _make_list_block(
                        current_list_type,
                        current_list_items,
                        shape,
                        slide_width_emu,
                        slide_height_emu,
                    )
                )
                current_list_items = []
                current_list_type = None

            block = {
                "type": block_type,
                "content": text,
            }
            if formatting:
                block["formatting"] = formatting
            block["position"] = _shape_position(shape, slide_width_emu, slide_height_emu)
            blocks.append(block)

    # Flush trailing list
    if current_list_items:
        blocks.append(
            _make_list_block(
                current_list_type,
                current_list_items,
                shape,
                slide_width_emu,
                slide_height_emu,
            )
        )

    return blocks


def _clean_bullet_text(text):
    """Remove leading bullet characters from text."""
    stripped = text.lstrip()
    if stripped and stripped[0] in BULLET_CHARS:
        stripped = stripped[1:].lstrip()
    # Also remove leading "- " or "* " patterns
    if stripped.startswith("- ") or stripped.startswith("* "):
        stripped = stripped[2:]
    return stripped


def _make_list_block(list_type, items, shape, slide_width_emu, slide_height_emu):
    """Create a list block dict."""
    block = {
        "type": list_type or "bullet_list",
        "content": "\n".join(items),
        "items": items,
        "position": _shape_position(shape, slide_width_emu, slide_height_emu),
    }
    return block


def _shape_position(shape, slide_width_emu, slide_height_emu):
    """Get shape position as percentage of slide dimensions."""
    try:
        return {
            "left_pct": _emu_to_pct(shape.left, slide_width_emu),
            "top_pct": _emu_to_pct(shape.top, slide_height_emu),
            "width_pct": _emu_to_pct(shape.width, slide_width_emu),
            "height_pct": _emu_to_pct(shape.height, slide_height_emu),
        }
    except Exception:
        return {"left_pct": 0, "top_pct": 0, "width_pct": 0, "height_pct": 0}


# ---------------------------------------------------------------------------
# Table extraction
# ---------------------------------------------------------------------------
def _extract_table(table, shape, slide_width_emu, slide_height_emu):
    """Extract table as markdown string."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip().replace("|", "\\|"))
        rows.append(cells)

    if not rows:
        return None

    # Build markdown table
    lines = []
    # Header row
    lines.append("| " + " | ".join(rows[0]) + " |")
    # Separator
    lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    # Data rows
    for row in rows[1:]:
        # Pad row if needed
        while len(row) < len(rows[0]):
            row.append("")
        lines.append("| " + " | ".join(row[: len(rows[0])]) + " |")

    md = "\n".join(lines)

    return {
        "type": "table",
        "content": md,
        "position": _shape_position(shape, slide_width_emu, slide_height_emu),
    }


# ---------------------------------------------------------------------------
# Chart extraction
# ---------------------------------------------------------------------------
def _extract_chart(chart, shape, slide_width_emu, slide_height_emu, slide_num, assets_dir):
    """
    Extract chart data: title + series data when possible.
    Falls back to saving chart as image placeholder.
    """
    result = {
        "type": "heading",
        "content": "",
        "position": _shape_position(shape, slide_width_emu, slide_height_emu),
    }

    try:
        chart_obj = chart
        title = ""
        if chart_obj.has_title and chart_obj.chart_title:
            try:
                title = chart_obj.chart_title.text_frame.text
            except Exception:
                title = "(chart)"

        series_data = []
        try:
            for series in chart_obj.series:
                sd = {"name": str(series.name) if series.name else "Series"}
                try:
                    sd["values"] = [
                        v for v in series.values if v is not None
                    ]
                except Exception:
                    sd["values"] = []
                series_data.append(sd)
        except Exception:
            pass

        # Try to get categories
        categories = []
        try:
            for plot in chart_obj.plots:
                for cat in plot.categories:
                    categories.append(str(cat))
                break
        except Exception:
            pass

        content_parts = []
        if title:
            content_parts.append(f"Chart: {title}")
        if categories:
            content_parts.append(f"Categories: {', '.join(categories)}")
        for sd in series_data:
            vals = ", ".join(str(v) for v in sd.get("values", []))
            content_parts.append(f"{sd['name']}: [{vals}]")

        result["content"] = "\n".join(content_parts) if content_parts else "(chart)"
        result["type"] = "body_text"  # Chart data presented as text

    except Exception as e:
        log.warning("  Chart extraction error: %s", e)
        result["content"] = "(chart - data extraction failed)"

    return result


# ---------------------------------------------------------------------------
# SmartArt extraction
# ---------------------------------------------------------------------------
def _extract_smartart(shape):
    """
    Attempt to extract text from SmartArt by parsing the shape's XML
    for dgm: namespace elements.
    """
    texts = []
    try:
        xml_str = shape._element.xml
        # Parse for text nodes within DrawingML / diagram namespace
        root = ET.fromstring(xml_str)

        # Search broadly for text content in any namespace
        for elem in root.iter():
            tag = elem.tag
            # Look for <a:t> text runs in any context
            if tag.endswith("}t") or tag == "t":
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
    except Exception as e:
        log.debug("  SmartArt XML parse error: %s", e)

    return texts


# ---------------------------------------------------------------------------
# Grouped shapes extraction (recursive)
# ---------------------------------------------------------------------------
def _extract_group_shapes(
    group_shape,
    slide_width_emu,
    slide_height_emu,
    slide_num,
    assets_dir,
    image_counter,
    extracted_images,
):
    """Recursively extract content from grouped shapes."""
    blocks = []
    images = []

    for shape in group_shape.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            sub_blocks, sub_images, image_counter = _extract_group_shapes(
                shape,
                slide_width_emu,
                slide_height_emu,
                slide_num,
                assets_dir,
                image_counter,
                extracted_images,
            )
            blocks.extend(sub_blocks)
            images.extend(sub_images)
        elif shape.has_text_frame:
            para_blocks = _extract_paragraphs(
                shape.text_frame,
                False,
                0,
                1,
                slide_width_emu,
                slide_height_emu,
                shape,
            )
            blocks.extend(para_blocks)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_id, image_counter = _extract_image(
                shape,
                slide_num,
                assets_dir,
                image_counter,
                slide_width_emu,
                slide_height_emu,
                extracted_images,
            )
            if img_id:
                images.append(img_id)

    return blocks, images, image_counter


# ---------------------------------------------------------------------------
# Image extraction
# ---------------------------------------------------------------------------
def _extract_image(
    shape,
    slide_num,
    assets_dir,
    image_counter,
    slide_width_emu,
    slide_height_emu,
    extracted_images,
):
    """Extract a single image from a picture shape. Returns (image_id, new_counter)."""
    try:
        image = shape.image
        image_bytes = image.blob
        image_ext = image.ext
        # Normalise extension
        if image_ext in ("wmf", "emf"):
            # Save as-is; downstream can convert
            pass
        image_counter += 1
        image_name = f"slide{slide_num}_img{image_counter}.{image_ext}"
        image_path = os.path.join(assets_dir, image_name)

        with open(image_path, "wb") as f:
            f.write(image_bytes)

        img_id = f"img_{slide_num}_{image_counter}"

        width_px = 0
        height_px = 0
        if Image and image_ext not in ("wmf", "emf"):
            try:
                pil_img = Image.open(BytesIO(image_bytes))
                width_px, height_px = pil_img.size
            except Exception:
                pass

        extracted_images.append(
            {
                "id": img_id,
                "path": f"assets/{image_name}",
                "source_slide": slide_num,
                "width": width_px,
                "height": height_px,
                "content_type": _guess_image_content_type(image_ext, width_px, height_px),
                "description": "",
                "regenerable": False,
            }
        )
        return img_id, image_counter
    except Exception as e:
        log.warning("  Image extraction error on slide %d: %s", slide_num, e)
        return None, image_counter


def _guess_image_content_type(ext, width, height):
    """Heuristic guess of image content type."""
    ext_lower = ext.lower() if ext else ""
    if ext_lower in ("svg",):
        return "diagram"
    if ext_lower in ("emf", "wmf"):
        return "diagram"
    if ext_lower == "ico":
        return "icon"
    # Small images are likely icons/logos
    if width and height and width < 100 and height < 100:
        return "icon"
    return "unknown"


# ---------------------------------------------------------------------------
# Slide layout detection
# ---------------------------------------------------------------------------
def _detect_slide_type(slide, content_blocks, images, has_table, has_chart):
    """
    Detect slide type:
    title, section_divider, content, two_column, image_heavy,
    code, chart, table, blank, closing
    """
    text_count = sum(
        1 for b in content_blocks if b.get("type") not in ("table",)
    )
    total_text = " ".join(b.get("content", "") for b in content_blocks)
    word_count = len(total_text.split())

    has_code = any(b.get("type") == "code_block" for b in content_blocks)
    has_heading = any(b.get("type") == "heading" for b in content_blocks)
    image_count = len(images)

    # Blank
    if text_count == 0 and image_count == 0:
        return "blank"

    # Chart
    if has_chart:
        return "chart"

    # Table
    if has_table:
        return "table"

    # Code-heavy
    if has_code and sum(1 for b in content_blocks if b["type"] == "code_block") >= text_count * 0.5:
        return "code"

    # Image-heavy: more images than text blocks, or large ratio
    if image_count > 0 and (image_count >= text_count or word_count < 30):
        return "image_heavy"

    # Title slide: has heading, few other elements, short text
    if has_heading and text_count <= 3 and word_count < 30:
        return "title"

    # Section divider: heading only, very few words
    if has_heading and text_count <= 2 and word_count < 15:
        return "section_divider"

    # Two-column detection: check if shapes span different halves
    left_blocks = [b for b in content_blocks if b.get("position", {}).get("left_pct", 0) < 40]
    right_blocks = [b for b in content_blocks if b.get("position", {}).get("left_pct", 0) > 50]
    if left_blocks and right_blocks and len(left_blocks) >= 2 and len(right_blocks) >= 2:
        return "two_column"

    # Closing slide heuristic
    lower_text = total_text.lower()
    closing_keywords = {"thank you", "thanks", "questions", "q&a", "contact", "the end"}
    if any(kw in lower_text for kw in closing_keywords) and word_count < 30:
        return "closing"

    return "content"


# ---------------------------------------------------------------------------
# Layout hints
# ---------------------------------------------------------------------------
def _compute_layout_hints(content_blocks, images):
    """Compute layout hints for a slide."""
    has_header = any(b.get("type") in ("heading", "subheading") for b in content_blocks)

    # Estimate column count from block positions
    left_positions = [
        b.get("position", {}).get("left_pct", 0) for b in content_blocks
    ]
    if left_positions:
        # Cluster into columns using a simple threshold
        sorted_pos = sorted(set(round(p / 10) * 10 for p in left_positions))
        column_count = max(1, len(sorted_pos))
        column_count = min(column_count, 3)  # Cap at 3
    else:
        column_count = 1

    # Dominant element
    type_counts = {}
    for b in content_blocks:
        t = b.get("type", "body_text")
        type_counts[t] = type_counts.get(t, 0) + 1
    if images and len(images) > sum(type_counts.values()):
        dominant = "image"
    elif type_counts:
        dominant = max(type_counts, key=type_counts.get)
    else:
        dominant = "image" if images else "empty"

    return {
        "has_header": has_header,
        "column_count": column_count,
        "dominant_element": dominant,
    }


# ---------------------------------------------------------------------------
# Slide rendering (reference PNGs)
# ---------------------------------------------------------------------------
def _render_slide_png(prs_path, slide_num, references_dir):
    """
    Render a slide as a reference PNG.
    Uses python-pptx slide dimensions + Pillow to create a placeholder
    reference image. For full-fidelity rendering, LibreOffice or
    comtypes (Windows) would be needed.

    Returns the relative path to the saved PNG, or None on failure.
    """
    png_name = f"slide_{slide_num:03d}.png"
    png_path = os.path.join(references_dir, png_name)

    # Attempt 1: Use comtypes on Windows to render via PowerPoint COM
    if sys.platform == "win32":
        try:
            return _render_slide_com(prs_path, slide_num, png_path)
        except Exception as e:
            log.debug("  COM rendering failed: %s", e)

    # Attempt 2: Use LibreOffice CLI
    try:
        return _render_slide_libreoffice(prs_path, slide_num, references_dir, png_name)
    except Exception:
        pass

    # Attempt 3: Create a placeholder thumbnail with Pillow
    if Image:
        try:
            img = Image.new("RGB", (960, 540), color=(240, 240, 240))
            img.save(png_path)
            return f"references/{png_name}"
        except Exception:
            pass

    return None


def _render_slide_com(prs_path, slide_num, png_path):
    """Render slide via PowerPoint COM automation (Windows only)."""
    import comtypes.client

    pptx_abs = os.path.abspath(prs_path)
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    try:
        presentation = powerpoint.Presentations.Open(pptx_abs, WithWindow=False)
        slide = presentation.Slides(slide_num)
        slide.Export(os.path.abspath(png_path), "PNG", 960, 540)
        presentation.Close()
    finally:
        # Don't quit PowerPoint if other instances may be open
        pass

    return f"references/{os.path.basename(png_path)}"


def _render_slide_libreoffice(prs_path, slide_num, references_dir, png_name):
    """Render all slides via LibreOffice (fallback)."""
    import subprocess
    import tempfile

    with tempfile.TemporaryDirectory() as tmpdir:
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "png",
                "--outdir",
                tmpdir,
                os.path.abspath(prs_path),
            ],
            capture_output=True,
            timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice failed: {result.stderr.decode()}")

        # LibreOffice produces one PNG per slide
        png_files = sorted(Path(tmpdir).glob("*.png"))
        if slide_num - 1 < len(png_files):
            import shutil

            dest = os.path.join(references_dir, png_name)
            shutil.copy2(str(png_files[slide_num - 1]), dest)
            return f"references/{png_name}"

    raise RuntimeError("No PNG produced by LibreOffice")


# ---------------------------------------------------------------------------
# OCR fallback
# ---------------------------------------------------------------------------
def _ocr_slide(prs_path, slide_num, references_dir, slide_width_emu, slide_height_emu):
    """
    For slides with minimal extractable text but images present,
    render the slide and run OCR to recover text.
    Returns a list of content_block dicts.
    """
    if not HAS_TESSERACT:
        log.info("  pytesseract not available; skipping OCR for slide %d", slide_num)
        return []

    # Try to get reference image
    png_name = f"slide_{slide_num:03d}.png"
    png_path = os.path.join(references_dir, png_name)

    if not os.path.exists(png_path):
        log.info("  No reference image for OCR on slide %d", slide_num)
        return []

    try:
        pil_img = Image.open(png_path)
        # Upscale for better OCR
        w, h = pil_img.size
        if w < 1920:
            scale = 1920 / w
            pil_img = pil_img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

        ocr_text = pytesseract.image_to_string(pil_img)
        if not ocr_text or len(ocr_text.strip()) < 5:
            return []

        # Split into blocks by double newlines
        raw_blocks = [b.strip() for b in ocr_text.split("\n\n") if b.strip()]
        blocks = []
        for rb in raw_blocks:
            blocks.append(
                {
                    "type": "body_text",
                    "content": rb,
                    "position": {
                        "left_pct": 0,
                        "top_pct": 0,
                        "width_pct": 100,
                        "height_pct": 100,
                    },
                }
            )
        return blocks
    except Exception as e:
        log.warning("  OCR failed for slide %d: %s", slide_num, e)
        return []


# ---------------------------------------------------------------------------
# Presentation metadata
# ---------------------------------------------------------------------------
def _extract_presentation_metadata(prs, file_path):
    """Extract presentation-level metadata."""
    core_props = prs.core_properties
    title = ""
    subtitle = ""
    author = ""

    # Try to get title from first slide
    if prs.slides:
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                if shape == first_slide.shapes.title:
                    title = shape.text.strip()
                elif not subtitle:
                    subtitle = shape.text.strip()

    # Core properties
    if core_props.author:
        author = core_props.author
    if core_props.title and not title:
        title = core_props.title

    # Simple language detection heuristic
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text + " "
            if len(all_text) > 2000:
                break

    detected_language = _detect_language(all_text)
    detected_subject = _detect_subject(all_text)

    return {
        "title": title,
        "subtitle": subtitle,
        "author": author,
        "slide_count": len(prs.slides),
        "source_type": "pptx",
        "source_file": os.path.basename(file_path),
        "detected_language": detected_language,
        "detected_subject": detected_subject,
    }


def _detect_language(text):
    """Simple heuristic language detection."""
    # Check for CJK characters
    cjk_count = sum(1 for c in text if "\u4e00" <= c <= "\u9fff")
    if cjk_count > len(text) * 0.1:
        return "zh"
    # Check for Korean
    kr_count = sum(1 for c in text if "\uac00" <= c <= "\ud7af")
    if kr_count > len(text) * 0.1:
        return "ko"
    # Check for Japanese hiragana/katakana
    jp_count = sum(
        1 for c in text if "\u3040" <= c <= "\u30ff" or "\u31f0" <= c <= "\u31ff"
    )
    if jp_count > len(text) * 0.05:
        return "ja"
    return "en"


def _detect_subject(text):
    """Simple heuristic subject detection based on keyword density."""
    lower = text.lower()
    subjects = {
        "science": [
            "hypothesis", "experiment", "data", "variable", "control",
            "observation", "theory", "molecule", "cell", "organism",
        ],
        "coding": [
            "function", "class", "variable", "algorithm", "api",
            "code", "programming", "software", "database", "deploy",
        ],
        "finance": [
            "revenue", "profit", "market", "investment", "stock",
            "portfolio", "gdp", "inflation", "dividend", "asset",
        ],
        "engineering": [
            "design", "system", "component", "specification", "testing",
            "requirement", "architecture", "prototype", "manufacturing",
        ],
        "medical": [
            "patient", "diagnosis", "treatment", "symptom", "clinical",
            "disease", "therapy", "medication", "surgery", "health",
        ],
        "humanities": [
            "culture", "history", "philosophy", "literature", "society",
            "art", "ethics", "narrative", "identity", "interpretation",
        ],
    }

    scores = {}
    for subject, keywords in subjects.items():
        score = sum(lower.count(kw) for kw in keywords)
        if score > 0:
            scores[subject] = score

    if scores:
        return max(scores, key=scores.get)
    return "general"


# ===========================================================================
# Main extraction
# ===========================================================================
def extract_pptx(file_path, output_dir="."):
    """
    Extract all content from a PowerPoint file.
    Produces slide-profile.json + assets/ + references/.
    """
    log.info("Opening %s", file_path)
    prs = Presentation(file_path)

    # Slide dimensions
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height

    # Create output directories
    assets_dir = os.path.join(output_dir, "assets")
    references_dir = os.path.join(output_dir, "references")
    os.makedirs(assets_dir, exist_ok=True)
    os.makedirs(references_dir, exist_ok=True)

    extracted_images = []
    slides_data = []

    # Presentation metadata
    pres_meta = _extract_presentation_metadata(prs, file_path)
    log.info(
        "Presentation: '%s' (%d slides)", pres_meta["title"], pres_meta["slide_count"]
    )

    for slide_num, slide in enumerate(prs.slides, start=1):
        log.info("Processing slide %d/%d...", slide_num, len(prs.slides))

        content_blocks = []
        slide_images = []
        image_counter = 0
        has_table = False
        has_chart = False
        slide_title = ""

        # Iterate shapes
        for shape_idx, shape in enumerate(slide.shapes):
            # ------- Text frames -------
            if shape.has_text_frame:
                is_title = shape == slide.shapes.title
                if is_title:
                    slide_title = shape.text.strip()

                para_blocks = _extract_paragraphs(
                    shape.text_frame,
                    is_title,
                    shape_idx,
                    len(slide.shapes),
                    slide_width_emu,
                    slide_height_emu,
                    shape,
                )
                content_blocks.extend(para_blocks)

            # ------- Tables -------
            if shape.has_table:
                has_table = True
                tbl_block = _extract_table(
                    shape.table, shape, slide_width_emu, slide_height_emu
                )
                if tbl_block:
                    content_blocks.append(tbl_block)

            # ------- Charts -------
            if shape.has_chart:
                has_chart = True
                chart_block = _extract_chart(
                    shape.chart,
                    shape,
                    slide_width_emu,
                    slide_height_emu,
                    slide_num,
                    assets_dir,
                )
                content_blocks.append(chart_block)

            # ------- Images (Picture shapes) -------
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_id, image_counter = _extract_image(
                    shape,
                    slide_num,
                    assets_dir,
                    image_counter,
                    slide_width_emu,
                    slide_height_emu,
                    extracted_images,
                )
                if img_id:
                    slide_images.append(img_id)

            # ------- Group shapes -------
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                grp_blocks, grp_images, image_counter = _extract_group_shapes(
                    shape,
                    slide_width_emu,
                    slide_height_emu,
                    slide_num,
                    assets_dir,
                    image_counter,
                    extracted_images,
                )
                content_blocks.extend(grp_blocks)
                slide_images.extend(grp_images)

            # ------- SmartArt / Diagrams -------
            # SmartArt shapes often have shape_type PLACEHOLDER or FREEFORM
            # Check XML for dgm: namespace
            try:
                xml_str = shape._element.xml
                if "dgm:" in xml_str or "diagram" in xml_str.lower():
                    texts = _extract_smartart(shape)
                    if texts:
                        # Combine as bullet list
                        content_blocks.append(
                            {
                                "type": "bullet_list",
                                "content": "\n".join(texts),
                                "items": texts,
                                "position": _shape_position(
                                    shape, slide_width_emu, slide_height_emu
                                ),
                            }
                        )
            except Exception:
                pass

        # ------- Speaker notes -------
        speaker_notes = ""
        if slide.has_notes_slide:
            try:
                speaker_notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

        # ------- Render reference PNG -------
        ref_path = _render_slide_png(file_path, slide_num, references_dir)
        if ref_path:
            log.info("  Reference image: %s", ref_path)

        # ------- OCR fallback -------
        total_text = " ".join(b.get("content", "") for b in content_blocks)
        if len(total_text.strip()) < 20 and (slide_images or len(slide.shapes) > 1):
            log.info("  Low text detected (%d chars); attempting OCR...", len(total_text.strip()))
            ocr_blocks = _ocr_slide(
                file_path, slide_num, references_dir, slide_width_emu, slide_height_emu
            )
            if ocr_blocks:
                log.info("  OCR recovered %d block(s)", len(ocr_blocks))
                content_blocks.extend(ocr_blocks)

        # ------- Detect slide type -------
        slide_type = _detect_slide_type(
            slide, content_blocks, slide_images, has_table, has_chart
        )

        # ------- Layout hints -------
        layout_hints = _compute_layout_hints(content_blocks, slide_images)

        # ------- Assemble slide data -------
        slide_data = {
            "number": slide_num,
            "title": slide_title,
            "slide_type": slide_type,
            "content_blocks": content_blocks,
            "speaker_notes": speaker_notes,
            "images": slide_images,
            "layout_hints": layout_hints,
        }
        slides_data.append(slide_data)

        # Console summary
        block_types = [b["type"] for b in content_blocks]
        log.info(
            "  Slide %d: '%s' [%s] — %d block(s), %d image(s)%s",
            slide_num,
            slide_title or "(untitled)",
            slide_type,
            len(content_blocks),
            len(slide_images),
            f", notes: {len(speaker_notes)} chars" if speaker_notes else "",
        )

    # ------- Build output -------
    output = {
        "presentation": pres_meta,
        "slides": slides_data,
        "extracted_images": extracted_images,
    }

    output_path = os.path.join(output_dir, "slide-profile.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output, f, indent=2, ensure_ascii=False)

    log.info("=" * 60)
    log.info("Extraction complete!")
    log.info("  Slides: %d", len(slides_data))
    log.info("  Images extracted: %d", len(extracted_images))
    log.info("  Output: %s", output_path)
    log.info("  Assets: %s", assets_dir)
    log.info("  References: %s", references_dir)

    return output


# ===========================================================================
# CLI
# ===========================================================================
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract-pptx.py <input.pptx> [output_dir]")
        print()
        print("Extracts all content from a PowerPoint file and produces:")
        print("  - slide-profile.json  (structured content data)")
        print("  - assets/             (extracted images)")
        print("  - references/         (rendered slide PNGs)")
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    if not os.path.isfile(input_file):
        log.error("File not found: %s", input_file)
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)
    extract_pptx(input_file, output_dir)
